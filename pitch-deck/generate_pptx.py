"""
SmartCareerAI Pitch Deck Generator - Light Theme with INR Currency
Generates a professional PowerPoint with light background and Indian Rupees
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# LIGHT color scheme
COLORS = {
    'primary': RGBColor(59, 130, 246),
    'secondary': RGBColor(139, 92, 246),
    'accent': RGBColor(6, 182, 212),
    'success': RGBColor(16, 185, 129),
    'warning': RGBColor(245, 158, 11),
    'danger': RGBColor(239, 68, 68),
    'bg': RGBColor(255, 255, 255),           # White background
    'card': RGBColor(248, 250, 252),         # Light gray card
    'text_primary': RGBColor(15, 23, 42),    # Dark text
    'text_secondary': RGBColor(71, 85, 105), # Medium gray text
    'text_muted': RGBColor(148, 163, 184),
    'border': RGBColor(226, 232, 240),
}

def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, text, left, top, width, height, font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color or COLORS['text_primary']
    p.alignment = align
    return txBox

def add_card(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['card']
    shape.line.color.rgb = COLORS['border']
    shape.line.width = Pt(1)
    return shape

def create_pitch_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ========== SLIDE 1: Title ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS['bg'])
    
    add_text(slide, "🚀 AI-Powered Career Platform", Inches(4), Inches(1.5), Inches(5.333), Inches(0.5), 
             font_size=14, color=COLORS['accent'], align=PP_ALIGN.CENTER)
    add_text(slide, "SmartCareerAI", Inches(1), Inches(2.2), Inches(11.333), Inches(1.2), 
             font_size=72, bold=True, color=COLORS['primary'], align=PP_ALIGN.CENTER)
    add_text(slide, "Empowering early-career students and professionals to land their dream jobs\nwith AI-driven insights, personalized coaching, and intelligent automation.",
             Inches(2), Inches(3.8), Inches(9.333), Inches(1), font_size=18, color=COLORS['text_secondary'], align=PP_ALIGN.CENTER)
    add_text(slide, "Pre-Seed  •  Prototype Stage  •  2026", 
             Inches(4), Inches(6.5), Inches(5.333), Inches(0.5), font_size=14, color=COLORS['text_muted'], align=PP_ALIGN.CENTER)
    
    # ========== SLIDE 2: Problem ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS['bg'])
    
    add_text(slide, "THE PROBLEM", Inches(5.5), Inches(0.5), Inches(2.333), Inches(0.4), 
             font_size=11, bold=True, color=COLORS['accent'], align=PP_ALIGN.CENTER)
    add_text(slide, "Job Hunting is Broken", Inches(1), Inches(1), Inches(11.333), Inches(0.8), 
             font_size=44, bold=True, align=PP_ALIGN.CENTER)
    
    problems = [
        ("75%", "of resumes never reach human eyes due to ATS rejection"),
        ("83%", "of students feel unprepared for interviews"),
        ("6+ Months", "average job search duration for fresh graduates"),
        ("Fragmented", "Tools for resumes, jobs, interviews are siloed"),
    ]
    
    for i, (stat, desc) in enumerate(problems):
        col, row = i % 2, i // 2
        left, top = Inches(1 + col * 5.5), Inches(2.2 + row * 2)
        add_card(slide, left, top, Inches(5), Inches(1.8))
        add_text(slide, stat, left + Inches(0.3), top + Inches(0.3), Inches(4.4), Inches(0.6),
                 font_size=28, bold=True, color=COLORS['primary'])
        add_text(slide, desc, left + Inches(0.3), top + Inches(1), Inches(4.4), Inches(0.6),
                 font_size=14, color=COLORS['text_secondary'])
    
    # ========== SLIDE 3: Solution ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS['bg'])
    
    add_text(slide, "THE SOLUTION", Inches(5.5), Inches(0.5), Inches(2.333), Inches(0.4), 
             font_size=11, bold=True, color=COLORS['accent'], align=PP_ALIGN.CENTER)
    add_text(slide, "One Platform. Complete Career Support.", Inches(1), Inches(1), Inches(11.333), Inches(0.8), 
             font_size=44, bold=True, align=PP_ALIGN.CENTER)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.666), Inches(2.8), Inches(2), Inches(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['primary']
    shape.line.fill.background()
    add_text(slide, "🚀", Inches(5.666), Inches(3.3), Inches(2), Inches(1), font_size=48, color=RGBColor(255,255,255), align=PP_ALIGN.CENTER)
    
    features = ["📄 Resume ATS Scoring", "🛤️ Skill Gap Analysis", "🏅 Skill Validation", 
                "💼 Job Aggregation", "🎥 AI Mock Interviews", "📧 Email Tracking"]
    positions = [(1, 2.5), (1, 4.5), (3.5, 5.5), (8, 5.5), (10, 4.5), (10, 2.5)]
    for text, (left, top) in zip(features, positions):
        add_text(slide, text, Inches(left), Inches(top), Inches(2.5), Inches(0.5),
                 font_size=14, align=PP_ALIGN.CENTER)
    
    # ========== SLIDE 4: Product Status ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS['bg'])
    
    add_text(slide, "PRODUCT", Inches(5.5), Inches(0.3), Inches(2.333), Inches(0.4), 
             font_size=11, bold=True, color=COLORS['accent'], align=PP_ALIGN.CENTER)
    add_text(slide, "Prototype Status", Inches(1), Inches(0.7), Inches(11.333), Inches(0.6), 
             font_size=40, bold=True, align=PP_ALIGN.CENTER)
    
    features = [
        ("📄 ATS Resume Scoring", "AI-powered resume analysis", "PROTOTYPE READY", COLORS['success']),
        ("🎤 AI Audio Interviews", "Voice-based practice", "PROTOTYPE READY", COLORS['success']),
        ("🛤️ Skill Gap Analysis", "Industry benchmarks", "PROTOTYPE READY", COLORS['success']),
        ("🔍 Job Aggregator", "Multi-platform search", "PROTOTYPE READY", COLORS['success']),
        ("🎥 AI Video Interviews", "Video mock interviews", "IN DEVELOPMENT", COLORS['warning']),
        ("📧 Email Tracking", "Interview email automation", "PLANNED", COLORS['text_muted']),
    ]
    
    for i, (title, desc, status, status_color) in enumerate(features):
        col, row = i % 3, i // 3
        left, top = Inches(0.5 + col * 4.2), Inches(1.5 + row * 2.8)
        add_card(slide, left, top, Inches(4), Inches(2.5))
        add_text(slide, status, left + Inches(0.15), top + Inches(0.15), Inches(1.8), Inches(0.3),
                 font_size=9, bold=True, color=status_color)
        add_text(slide, title, left + Inches(0.15), top + Inches(0.5), Inches(3.7), Inches(0.5),
                 font_size=16, bold=True)
        add_text(slide, desc, left + Inches(0.15), top + Inches(1.1), Inches(3.7), Inches(1.2),
                 font_size=12, color=COLORS['text_secondary'])
    
    # ========== SLIDE 5: Traction ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS['bg'])
    
    add_text(slide, "TRACTION", Inches(5.5), Inches(0.5), Inches(2.333), Inches(0.4), 
             font_size=11, bold=True, color=COLORS['accent'], align=PP_ALIGN.CENTER)
    add_text(slide, "Validation & Early Signals", Inches(1), Inches(1), Inches(11.333), Inches(0.8), 
             font_size=44, bold=True, align=PP_ALIGN.CENTER)
    
    validations = [
        ("👥 Validation Interviews", ["Interviews: [X] students", "Pain confirmed: [X]%", "Willingness to pay: [X]%"]),
        ("📊 Survey Results", ["Respondents: [X] students", "Top pain: ATS rejection", "Interest: [X]%"]),
        ("🧪 Prototype Testing", ["Testers: [X] users", "Task completion: [X]%", "NPS Score: [X]"]),
    ]
    
    for i, (title, items) in enumerate(validations):
        left = Inches(0.5 + i * 4.2)
        add_card(slide, left, Inches(2.2), Inches(4), Inches(3.5))
        add_text(slide, title, left + Inches(0.2), Inches(2.4), Inches(3.6), Inches(0.5),
                 font_size=16, bold=True)
        for j, item in enumerate(items):
            add_text(slide, item, left + Inches(0.2), Inches(3 + j * 0.5), Inches(3.6), Inches(0.4),
                     font_size=12, color=COLORS['text_secondary'])
    
    add_text(slide, "⚠️ PROTOTYPE STAGE - metrics represent validation research",
             Inches(1), Inches(6.2), Inches(11.333), Inches(0.5), font_size=12, color=COLORS['warning'], align=PP_ALIGN.CENTER)
    
    # ========== SLIDE 6: Business Model ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS['bg'])
    
    add_text(slide, "BUSINESS MODEL", Inches(5), Inches(0.5), Inches(3.333), Inches(0.4), 
             font_size=11, bold=True, color=COLORS['accent'], align=PP_ALIGN.CENTER)
    add_text(slide, "How We Make Money", Inches(1), Inches(1), Inches(11.333), Inches(0.8), 
             font_size=44, bold=True, align=PP_ALIGN.CENTER)
    
    tiers = [("Free", "₹0/mo", "Acquisition"), ("Pro", "₹499/mo", "Core Revenue"), ("Enterprise", "Custom", "B2B")]
    for i, (name, price, purpose) in enumerate(tiers):
        left = Inches(0.5 + i * 2.8)
        add_card(slide, left, Inches(2), Inches(2.5), Inches(2.5))
        add_text(slide, name, left, Inches(2.2), Inches(2.5), Inches(0.4), font_size=16, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, price, left, Inches(2.6), Inches(2.5), Inches(0.5), font_size=24, bold=True, color=COLORS['primary'], align=PP_ALIGN.CENTER)
        add_text(slide, purpose, left, Inches(4), Inches(2.5), Inches(0.3), font_size=10, color=COLORS['text_muted'], align=PP_ALIGN.CENTER)
    
    add_card(slide, Inches(8.5), Inches(2), Inches(4.333), Inches(2.5))
    add_text(slide, "Unit Economics (Projected)", Inches(8.7), Inches(2.2), Inches(4), Inches(0.4), font_size=14, bold=True)
    add_text(slide, "Target LTV: ₹18,000\nTarget CAC: ₹2,000-3,500\nLTV:CAC: 5-9x", 
             Inches(8.7), Inches(2.7), Inches(4), Inches(1.5), font_size=12, color=COLORS['text_secondary'])
    
    add_card(slide, Inches(0.5), Inches(5), Inches(12.333), Inches(2))
    add_text(slide, "Acquisition Strategy: University partnerships • Organic content • Student ambassadors • Referrals",
             Inches(0.7), Inches(5.5), Inches(12), Inches(1), font_size=14, color=COLORS['text_secondary'])
    
    # ========== SLIDE 7: Market ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS['bg'])
    
    add_text(slide, "MARKET OPPORTUNITY", Inches(4.5), Inches(0.5), Inches(4.333), Inches(0.4), 
             font_size=11, bold=True, color=COLORS['accent'], align=PP_ALIGN.CENTER)
    add_text(slide, "Massive TAM", Inches(1), Inches(1), Inches(11.333), Inches(0.8), 
             font_size=44, bold=True, align=PP_ALIGN.CENTER)
    
    markets = [("₹4.2L Cr", "TAM", "Global HR Tech"), ("₹1L Cr", "SAM", "Career Services"), ("₹6,700 Cr", "SOM", "Student Career")]
    for i, (amount, label, desc) in enumerate(markets):
        left = Inches(1 + i * 4)
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, Inches(2.2), Inches(3.5), Inches(3.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS['primary'] if i == 2 else COLORS['card']
        shape.line.color.rgb = COLORS['primary']
        add_text(slide, amount, left, Inches(3.2), Inches(3.5), Inches(0.6), font_size=28, bold=True, 
                 color=RGBColor(255,255,255) if i == 2 else COLORS['text_primary'], align=PP_ALIGN.CENTER)
        add_text(slide, label, left, Inches(3.8), Inches(3.5), Inches(0.4), font_size=14, bold=True,
                 color=RGBColor(255,255,255) if i == 2 else COLORS['text_primary'], align=PP_ALIGN.CENTER)
        add_text(slide, desc, left, Inches(4.2), Inches(3.5), Inches(0.4), font_size=12,
                 color=RGBColor(255,255,255) if i == 2 else COLORS['text_secondary'], align=PP_ALIGN.CENTER)
    
    stats = [("40M+", "Annual graduates"), ("23%", "AI recruitment CAGR"), ("₹40,000", "Avg spend/student")]
    for i, (val, lbl) in enumerate(stats):
        add_card(slide, Inches(1 + i * 4), Inches(6), Inches(3.5), Inches(1.2))
        add_text(slide, val, Inches(1 + i * 4), Inches(6.1), Inches(3.5), Inches(0.5), font_size=20, bold=True, color=COLORS['primary'], align=PP_ALIGN.CENTER)
        add_text(slide, lbl, Inches(1 + i * 4), Inches(6.6), Inches(3.5), Inches(0.4), font_size=11, color=COLORS['text_secondary'], align=PP_ALIGN.CENTER)
    
    # ========== SLIDE 8: Competition ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS['bg'])
    
    add_text(slide, "COMPETITIVE LANDSCAPE", Inches(4.5), Inches(0.5), Inches(4.333), Inches(0.4), 
             font_size=11, bold=True, color=COLORS['accent'], align=PP_ALIGN.CENTER)
    add_text(slide, "Why We Win", Inches(1), Inches(1), Inches(11.333), Inches(0.8), 
             font_size=44, bold=True, align=PP_ALIGN.CENTER)
    
    headers = ["Features", "🚀 SmartCareerAI", "Resume.io", "Pramp", "LinkedIn"]
    for i, h in enumerate(headers):
        add_text(slide, h, Inches(0.5 + i * 2.5), Inches(2), Inches(2.3), Inches(0.4),
                 font_size=12, bold=True, color=COLORS['primary'] if i == 1 else COLORS['text_primary'], align=PP_ALIGN.CENTER)
    
    rows = [("ATS Resume Scoring", "✓", "✓", "✗", "✗"), ("AI Mock Interviews", "✓", "✗", "✓", "✗"),
            ("Skill Gap Analysis", "✓", "✗", "✗", "~"), ("Job Aggregation", "✓", "✗", "✗", "✓"), ("Unified Platform", "✓", "✗", "✗", "✗")]
    for r, row in enumerate(rows):
        for c, cell in enumerate(row):
            color = COLORS['success'] if cell == "✓" else (COLORS['warning'] if cell == "~" else COLORS['text_muted'])
            add_text(slide, cell, Inches(0.5 + c * 2.5), Inches(2.5 + r * 0.55), Inches(2.3), Inches(0.4),
                     font_size=12, color=color if c > 0 else COLORS['text_primary'], align=PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT)
    
    add_card(slide, Inches(0.5), Inches(5.3), Inches(12.333), Inches(1.9))
    add_text(slide, "🛡️ Unfair Advantage: All-in-One Platform • AI-Native • Student-First Pricing", 
             Inches(0.7), Inches(5.6), Inches(11.5), Inches(1), font_size=14, bold=True)
    
    # ========== SLIDE 9: GTM ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS['bg'])
    
    add_text(slide, "GO-TO-MARKET", Inches(5), Inches(0.5), Inches(3.333), Inches(0.4), 
             font_size=11, bold=True, color=COLORS['accent'], align=PP_ALIGN.CENTER)
    add_text(slide, "Path to First 1000 Users", Inches(1), Inches(1), Inches(11.333), Inches(0.8), 
             font_size=44, bold=True, align=PP_ALIGN.CENTER)
    
    channels = [
        ("📢 Acquisition", "LinkedIn • YouTube • Reddit • SEO"),
        ("🤝 Partnerships", "Universities • Bootcamps • Career coaches"),
        ("🚀 Launch", "Product Hunt • Hacker News • Influencers"),
    ]
    for i, (title, items) in enumerate(channels):
        left = Inches(0.5 + i * 4.2)
        add_card(slide, left, Inches(2), Inches(4), Inches(2.5))
        add_text(slide, title, left + Inches(0.2), Inches(2.2), Inches(3.6), Inches(0.5), font_size=16, bold=True)
        add_text(slide, items, left + Inches(0.2), Inches(2.8), Inches(3.6), Inches(1.5), font_size=12, color=COLORS['text_secondary'])
    
    add_card(slide, Inches(0.5), Inches(5), Inches(12.333), Inches(2.2))
    add_text(slide, "First 1000: Beta testers → Iterate → Public launch → Referrals", 
             Inches(0.7), Inches(5.5), Inches(11.5), Inches(1), font_size=16, align=PP_ALIGN.CENTER)
    
    # ========== SLIDE 10: Team ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS['bg'])
    
    add_text(slide, "TEAM", Inches(5.5), Inches(0.5), Inches(2.333), Inches(0.4), 
             font_size=11, bold=True, color=COLORS['accent'], align=PP_ALIGN.CENTER)
    add_text(slide, "The People Behind It", Inches(1), Inches(1), Inches(11.333), Inches(0.8), 
             font_size=44, bold=True, align=PP_ALIGN.CENTER)
    
    founders = [("[Founder Name]", "CEO & Co-Founder", "[Your background]"),
                ("[Co-Founder Name]", "CTO & Co-Founder", "[Your background]")]
    for i, (name, role, bio) in enumerate(founders):
        left = Inches(2.5 + i * 4.5)
        add_card(slide, left, Inches(2), Inches(4), Inches(3))
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(1.5), Inches(2.3), Inches(1), Inches(1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS['primary']
        shape.line.fill.background()
        add_text(slide, name, left, Inches(3.5), Inches(4), Inches(0.4), font_size=16, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, role, left, Inches(3.9), Inches(4), Inches(0.4), font_size=12, color=COLORS['primary'], align=PP_ALIGN.CENTER)
        add_text(slide, bio, left + Inches(0.2), Inches(4.4), Inches(3.6), Inches(0.5), font_size=11, color=COLORS['text_secondary'], align=PP_ALIGN.CENTER)
    
    add_card(slide, Inches(0.5), Inches(5.3), Inches(12.333), Inches(1.9))
    add_text(slide, "Why Us: Lived the problem • Technical depth • Domain access", 
             Inches(0.7), Inches(5.8), Inches(11.5), Inches(0.8), font_size=14, bold=True, align=PP_ALIGN.CENTER)
    
    # ========== SLIDE 11: Financials ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS['bg'])
    
    add_text(slide, "FINANCIALS", Inches(5.5), Inches(0.5), Inches(2.333), Inches(0.4), 
             font_size=11, bold=True, color=COLORS['accent'], align=PP_ALIGN.CENTER)
    add_text(slide, "12-24 Month Projections", Inches(1), Inches(1), Inches(11.333), Inches(0.8), 
             font_size=44, bold=True, align=PP_ALIGN.CENTER)
    
    add_card(slide, Inches(0.5), Inches(2), Inches(5), Inches(4))
    add_text(slide, "📊 Revenue Projections", Inches(0.7), Inches(2.2), Inches(4.6), Inches(0.4), font_size=14, bold=True)
    proj_data = [("Metric", "M6", "M12", "M24"), ("Users", "500", "2K", "10K"), ("Paid", "50", "300", "2K"),
                 ("MRR", "₹79K", "₹4.75L", "₹31.7L"), ("ARR", "₹9.5L", "₹57L", "₹3.8Cr")]
    for r, row in enumerate(proj_data):
        for c, cell in enumerate(row):
            bold = r == 0 or r == 4
            color = COLORS['primary'] if r == 4 else (COLORS['text_muted'] if r == 0 else COLORS['text_primary'])
            add_text(slide, cell, Inches(0.7 + c * 1.1), Inches(2.7 + r * 0.5), Inches(1), Inches(0.4), font_size=11, bold=bold, color=color)
    
    add_card(slide, Inches(5.7), Inches(2), Inches(3.5), Inches(4))
    add_text(slide, "🔢 Assumptions", Inches(5.9), Inches(2.2), Inches(3.1), Inches(0.4), font_size=14, bold=True)
    assumptions = ["10% conversion", "₹499/mo avg", "5% churn", "₹2K CAC", "50% MoM growth"]
    for i, a in enumerate(assumptions):
        add_text(slide, f"• {a}", Inches(5.9), Inches(2.7 + i * 0.45), Inches(3.1), Inches(0.4), font_size=11, color=COLORS['text_secondary'])
    
    add_card(slide, Inches(9.4), Inches(2), Inches(3.433), Inches(4))
    add_text(slide, "🔥 Burn & Runway", Inches(9.6), Inches(2.2), Inches(3), Inches(0.4), font_size=14, bold=True)
    burns = [("Monthly Burn", "₹12-17L"), ("Primary Costs", "Eng, Cloud, AI"), ("Break-even", "Month 18-24")]
    for i, (lbl, val) in enumerate(burns):
        add_text(slide, lbl, Inches(9.6), Inches(2.8 + i * 0.8), Inches(3), Inches(0.3), font_size=10, color=COLORS['text_secondary'])
        add_text(slide, val, Inches(9.6), Inches(3.1 + i * 0.8), Inches(3), Inches(0.4), font_size=14, bold=True)
    
    # ========== SLIDE 12: The Ask ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS['bg'])
    
    add_text(slide, "THE ASK", Inches(5.5), Inches(0.5), Inches(2.333), Inches(0.4), 
             font_size=11, bold=True, color=COLORS['accent'], align=PP_ALIGN.CENTER)
    add_text(slide, "Raising ₹2 Crore Pre-Seed", Inches(1), Inches(1), Inches(11.333), Inches(0.8), 
             font_size=44, bold=True, align=PP_ALIGN.CENTER)
    
    add_text(slide, "Use of Funds", Inches(0.5), Inches(2), Inches(4), Inches(0.5), font_size=18, bold=True)
    funds = [("50%", "Product Development"), ("25%", "Go-to-Market"), ("15%", "Infrastructure"), ("10%", "Reserve")]
    for i, (pct, name) in enumerate(funds):
        add_text(slide, f"{pct} {name}", Inches(0.5), Inches(2.6 + i * 0.7), Inches(6), Inches(0.4), font_size=14)
    
    add_text(slide, "12-Month Milestones", Inches(7), Inches(2), Inches(6), Inches(0.5), font_size=18, bold=True)
    milestones = [("2,000", "Users"), ("₹4L", "MRR"), ("5", "Universities"), ("Seed", "Ready")]
    for i, (val, lbl) in enumerate(milestones):
        col, row = i % 2, i // 2
        add_card(slide, Inches(7 + col * 3), Inches(2.6 + row * 1.8), Inches(2.8), Inches(1.6))
        add_text(slide, val, Inches(7 + col * 3), Inches(2.8 + row * 1.8), Inches(2.8), Inches(0.5), font_size=24, bold=True, color=COLORS['primary'], align=PP_ALIGN.CENTER)
        add_text(slide, lbl, Inches(7 + col * 3), Inches(3.3 + row * 1.8), Inches(2.8), Inches(0.4), font_size=12, color=COLORS['text_secondary'], align=PP_ALIGN.CENTER)
    
    add_text(slide, "💰 Runway: 12-15 months", Inches(7), Inches(6.3), Inches(6), Inches(0.4), font_size=14, color=COLORS['success'], align=PP_ALIGN.CENTER)
    
    # ========== SLIDE 13: Vision & Risks ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS['bg'])
    
    add_text(slide, "VISION & RISKS", Inches(5), Inches(0.5), Inches(3.333), Inches(0.4), 
             font_size=11, bold=True, color=COLORS['accent'], align=PP_ALIGN.CENTER)
    add_text(slide, "Long-term Roadmap", Inches(1), Inches(1), Inches(11.333), Inches(0.8), 
             font_size=44, bold=True, align=PP_ALIGN.CENTER)
    
    add_card(slide, Inches(0.5), Inches(2), Inches(6), Inches(4.7))
    add_text(slide, "🔭 Vision", Inches(0.7), Inches(2.2), Inches(5.6), Inches(0.4), font_size=16, bold=True)
    visions = [("Year 1", "MVP → Product-Market Fit"), ("Year 2", "Scale to 50K users"), ("Year 3+", "International expansion")]
    for i, (yr, goal) in enumerate(visions):
        add_text(slide, f"{yr}: {goal}", Inches(0.7), Inches(2.8 + i * 0.6), Inches(5.6), Inches(0.5), font_size=12, color=COLORS['text_secondary'])
    add_text(slide, "Goal: Default career platform for every student globally", 
             Inches(0.7), Inches(4.8), Inches(5.6), Inches(0.8), font_size=12, bold=True, color=COLORS['primary'])
    
    add_card(slide, Inches(6.833), Inches(2), Inches(6), Inches(4.7))
    add_text(slide, "⚠️ Risks & Mitigation", Inches(7.033), Inches(2.2), Inches(5.6), Inches(0.4), font_size=16, bold=True)
    risks = [("Competition", "Move fast, unified experience"), ("AI costs", "Optimize prompts, caching"), ("CAC", "Organic-first, community")]
    for i, (risk, mitigation) in enumerate(risks):
        add_text(slide, f"• {risk}: {mitigation}", Inches(7.033), Inches(2.8 + i * 0.7), Inches(5.6), Inches(0.6), font_size=12, color=COLORS['text_secondary'])
    
    # ========== SLIDE 14: Contact ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS['bg'])
    
    add_text(slide, "Let's Build the Future of\nCareer Success", Inches(1), Inches(1.5), Inches(11.333), Inches(1.5), 
             font_size=48, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Join us in helping millions of students land their dream jobs.",
             Inches(1), Inches(3.3), Inches(11.333), Inches(0.8), font_size=18, color=COLORS['text_secondary'], align=PP_ALIGN.CENTER)
    
    contacts = ["📧 [your@email.com]", "🌐 [website]", "💼 [LinkedIn]"]
    for i, c in enumerate(contacts):
        add_card(slide, Inches(1.5 + i * 3.5), Inches(4.5), Inches(3.2), Inches(0.6))
        add_text(slide, c, Inches(1.5 + i * 3.5), Inches(4.55), Inches(3.2), Inches(0.5), font_size=14, align=PP_ALIGN.CENTER)
    
    add_text(slide, "🚀 SmartCareerAI", Inches(1), Inches(5.8), Inches(11.333), Inches(0.5), font_size=32, bold=True, color=COLORS['primary'], align=PP_ALIGN.CENTER)
    add_text(slide, "AI-Powered Career Success Platform", Inches(1), Inches(6.4), Inches(11.333), Inches(0.4), font_size=14, color=COLORS['text_secondary'], align=PP_ALIGN.CENTER)
    
    # Save
    output_path = os.path.join(os.path.dirname(__file__), "SmartCareerAI_Pitch_Deck_Light_INR.pptx")
    prs.save(output_path)
    print(f"✅ Light theme pitch deck saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    create_pitch_deck()
